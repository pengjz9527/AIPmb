"""核心视觉组件：标题栏、卡片、表格、代码块、架构图、流程图等 — 居中平衡设计"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

from .theme import *

# ============================================================
# 标题栏：标题左 + Logo右
# ============================================================
def add_header(slide, title, section_num=""):
    """添加页面标题栏：左边accent竖线+标题，右边logo"""
    line_left = MARGIN_EMU
    line_h = HEADER_H_EMU - px(10)
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, line_left, Emu(MARGIN_EMU + px(5)), px(5), line_h)
    set_shape_fill(line_shape, ACCENT)
    line_shape.line.fill.background()

    title_left = line_left + px(16)
    title_w = SLIDE_W_EMU - 2 * MARGIN_EMU - LOGO_W_EMU - px(16)
    full_title = f"{section_num} {title}" if section_num else title
    add_textbox(slide, title_left, Emu(MARGIN_EMU), title_w, HEADER_H_EMU,
                full_title, font_size=H2_SZ, color=TEXT_DARK, bold=True)

    logo_left = SLIDE_W_EMU - MARGIN_EMU - LOGO_W_EMU
    logo_top = Emu(MARGIN_EMU + (HEADER_H - LOGO_H) // 2 * PX_TO_EMU)
    try:
        slide.shapes.add_picture(LOGO_PATH_ABS, logo_left, logo_top, LOGO_W_EMU, LOGO_H_EMU)
    except FileNotFoundError:
        pass

    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_EMU, Emu(MARGIN_EMU + HEADER_H_EMU),
        CONTENT_W_EMU, px(1))
    set_shape_fill(div, GRAY_200)
    div.line.fill.background()


# ============================================================
# 封面页
# ============================================================
def add_cover_slide(slide, title, subtitle, date_text, footer_text):
    """深色背景封面/感谢页"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK

    # Logo 居中
    logo_w = px(220)
    logo_h = px(89)
    try:
        slide.shapes.add_picture(LOGO_PATH_ABS, center_x(logo_w), px(100), logo_w, logo_h)
    except FileNotFoundError:
        pass

    # 主标题
    add_textbox(slide, MARGIN_EMU, px(230), CONTENT_W_EMU, px(90),
                title, H1_SZ, WHITE, True, PP_ALIGN.CENTER)

    # 装饰线
    bar_w = px(100)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x(bar_w), px(340), bar_w, px(5))
    set_shape_fill(bar, ACCENT)
    bar.line.fill.background()

    # 副标题
    add_textbox(slide, MARGIN_EMU, px(370), CONTENT_W_EMU, px(60),
                subtitle, Pt(26), RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)

    # 日期
    add_textbox(slide, MARGIN_EMU, px(460), CONTENT_W_EMU, px(40),
                date_text, Pt(16), RGBColor(0x88, 0x88, 0x88), False, PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, MARGIN_EMU, px(730), CONTENT_W_EMU, px(32),
                footer_text, Pt(12), RGBColor(0x88, 0x88, 0x88), False, PP_ALIGN.CENTER)


# ============================================================
# 章节分隔页
# ============================================================
def add_divider_slide(slide, num, title):
    """蓝色背景章节分隔页"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT

    add_textbox(slide, MARGIN_EMU, px(180), CONTENT_W_EMU, px(200),
                num, Pt(100), WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN_EMU, px(440), CONTENT_W_EMU, px(80),
                title, Pt(42), WHITE, True, PP_ALIGN.CENTER)


# ============================================================
# 卡片网格 — 居中平衡
# ============================================================
def add_card_grid(slide, cards, cols=2, card_colors=None, top_emu=None):
    """通用卡片网格，内容居中分布"""
    if top_emu is None:
        top_emu = BODY_TOP_EMU
    if card_colors is None:
        card_colors = [ACCENT, BLUE, GREEN, PURPLE, ORANGE, ACCENT]

    gap = px(20)
    n_rows = (len(cards) + cols - 1) // cols
    card_w = (CONTENT_W_EMU - (cols - 1) * gap) // cols
    avail_h = SLIDE_H_EMU - top_emu - MARGIN_EMU
    card_h = min(avail_h // n_rows - gap // n_rows, px(160) if n_rows == 2 else px(280))

    # 整个网格居中
    total_w = cols * card_w + (cols - 1) * gap
    grid_left = center_x(total_w)

    for i, card in enumerate(cards):
        row = i // cols
        col = i % cols
        left = grid_left + col * (card_w + gap)
        top = top_emu + row * (card_h + gap)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        set_shape_fill(shape, WHITE)
        set_shape_border(shape, GRAY_200, 0.5)

        color = card_colors[i % len(card_colors)]
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, card_w, px(4))
        set_shape_fill(accent_line, color)
        accent_line.line.fill.background()

        # 卡片内容 — 居中排列
        inner_pad = px(20)
        inner_l = left + inner_pad
        inner_t = top + px(16)
        inner_w = card_w - 2 * inner_pad
        inner_h = card_h - px(36)

        tb = add_textbox(slide, inner_l, inner_t, inner_w, inner_h)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if card.get("icon"):
            run_icon = p.add_run()
            run_icon.text = card["icon"] + "\n"
            run_icon.font.size = Pt(28)
            run_icon.font.name = FONT_FAMILY
        run_title = p.add_run()
        run_title.text = card["title"]
        run_title.font.size = H4_SZ
        run_title.font.bold = True
        run_title.font.color.rgb = TEXT_DARK
        run_title.font.name = FONT_FAMILY

        if card.get("desc"):
            p2 = tb.text_frame.add_paragraph()
            p2.text = card["desc"]
            p2.font.size = SMALL_SZ
            p2.font.color.rgb = TEXT_GRAY
            p2.font.name = FONT_FAMILY
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(8)
            p2.space_after = Pt(0)


# ============================================================
# 架构图 — 居中
# ============================================================
def add_arch_diagram(slide, layers, top_emu=None):
    """纵向分层架构图，居中"""
    if top_emu is None:
        top_emu = BODY_TOP_EMU

    arch_w = px(1100)
    arch_left = center_x(arch_w)
    label_h = px(18)
    layer_h_base = px(50)
    gap = px(4)
    arrow_h = px(20)

    n_layers = len(layers)
    total_avail = SLIDE_H_EMU - top_emu - MARGIN_EMU
    fixed_h = n_layers * (label_h + layer_h_base + gap) + (n_layers - 1) * arrow_h
    scale = 1.0
    if fixed_h > total_avail:
        scale = total_avail / fixed_h
    layer_h = int(layer_h_base * scale) if scale < 1.0 else layer_h_base
    current_y = top_emu

    for li, layer in enumerate(layers):
        # 层标签
        add_textbox(slide, arch_left, current_y, arch_w, label_h,
                    layer["label"], CAPTION_SZ, TEXT_GRAY, True, PP_ALIGN.CENTER)
        current_y += int(px(16) * min(scale, 1)) + gap

        # 层背景
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, arch_left, current_y, arch_w, layer_h)
        set_shape_fill(bg_shape, layer.get("bg_color", GRAY_100))
        bg_shape.line.fill.background()

        boxes = layer.get("boxes", [])
        if boxes:
            box_gap = px(10)
            box_pad = px(12)
            total_box_w = arch_w - 2 * box_pad
            box_w = (total_box_w - (len(boxes) - 1) * box_gap) // len(boxes)
            box_h = layer_h - px(8)

            for bi, box in enumerate(boxes):
                b_left = arch_left + box_pad + bi * (box_w + box_gap)
                b_top = current_y + px(4)

                b_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, b_left, b_top, box_w, box_h)
                set_shape_fill(b_shape, box.get("color", ACCENT_LIGHT))
                b_shape.line.fill.background()

                b_tb = slide.shapes.add_textbox(b_left, b_top, box_w, box_h)
                b_tb.word_wrap = True
                tf = b_tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = box.get("text", "")
                run.font.size = Pt(12) if scale > 0.9 else Pt(10)
                run.font.bold = True
                run.font.color.rgb = box.get("text_color", TEXT_DARK)
                run.font.name = FONT_FAMILY

                if box.get("sub"):
                    p2 = tf.add_paragraph()
                    p2.alignment = PP_ALIGN.CENTER
                    run2 = p2.add_run()
                    run2.text = box["sub"]
                    run2.font.size = Pt(9) if scale > 0.9 else Pt(7)
                    run2.font.color.rgb = TEXT_GRAY
                    run2.font.name = FONT_FAMILY

        current_y += layer_h + gap

        if li < len(layers) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                SLIDE_CENTER_EMU - px(14), current_y, px(28), int(arrow_h * min(scale, 1)))
            set_shape_fill(arrow, ACCENT)
            arrow.line.fill.background()
            current_y += int(arrow_h * min(scale, 1)) + gap


# ============================================================
# 表格 — 居中
# ============================================================
def add_table(slide, headers, rows, col_widths_emu, top_emu=None):
    """通用表格，居中"""
    if top_emu is None:
        top_emu = BODY_TOP_EMU

    table_width = sum(col_widths_emu)
    left = center_x(table_width)
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    row_h = px(38)
    total_h = n_rows * row_h

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top_emu, table_width, total_h)
    tbl = tbl_shape.table

    for ci, w in enumerate(col_widths_emu):
        tbl.columns[ci].width = w

    # 表头
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK

    # 数据行
    for ri, row in enumerate(rows):
        for ci, text in enumerate(row.get("cells", [])):
            cell = tbl.cell(ri + 1, ci)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_DARK
                p.font.name = FONT_FAMILY
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_50
        hl_col = row.get("highlight_col")
        if hl_col is not None:
            cell = tbl.cell(ri + 1, hl_col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_LIGHT


# ============================================================
# 代码块
# ============================================================
def add_code_block(slide, code, left_emu, top_emu, width_emu, height_emu, lang="Python"):
    """深色背景代码块"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_emu, top_emu, width_emu, height_emu)
    set_shape_fill(shape, CODE_BG)
    shape.line.fill.background()

    # 语言标签
    lang_w = px(64)
    lang_h = px(20)
    lang_badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_emu + px(10), top_emu + px(8), lang_w, lang_h)
    set_shape_fill(lang_badge, DARK_SECONDARY)
    lang_badge.line.fill.background()
    add_textbox(slide, left_emu + px(10), top_emu + px(8), lang_w, lang_h,
                lang, Pt(9), TEXT_GRAY, False, PP_ALIGN.CENTER)

    # 代码文本
    code_left = left_emu + px(16)
    code_top = top_emu + px(32)
    code_w = width_emu - px(32)
    code_h = height_emu - px(40)

    tb = slide.shapes.add_textbox(code_left, code_top, code_w, code_h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True

    lines = code.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = CODE_SZ
        p.font.color.rgb = CODE_TEXT
        p.font.name = CODE_FONT
        p.space_before = Pt(0)
        p.space_after = Pt(0)


# ============================================================
# 评分条 — 居中
# ============================================================
def add_score_bars(slide, scores, top_emu=None):
    """评分条对比，居中"""
    if top_emu is None:
        top_emu = BODY_TOP_EMU

    label_w = px(160)
    bar_w = px(700)
    val_w = px(60)
    row_h = px(50)
    gap = px(10)

    total_w = label_w + px(20) + bar_w + val_w * 2 + px(20)
    start_x = center_x(total_w)

    for i, sc in enumerate(scores):
        y = top_emu + i * (row_h + gap)

        add_textbox(slide, start_x, y, label_w, row_h,
                    sc["label"], BODY_SZ, TEXT_DARK, True, PP_ALIGN.RIGHT)

        bg_left = start_x + label_w + px(20)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bg_left, y + px(16), bar_w, px(14))
        set_shape_fill(bg_shape, GRAY_200)
        bg_shape.line.fill.background()

        aipmb_w = int(bar_w * sc["aipmb_pct"] / 100)
        if aipmb_w > 0:
            fg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, bg_left, y + px(16), aipmb_w, px(14))
            set_shape_fill(fg, ACCENT)
            fg.line.fill.background()

        add_textbox(slide, bg_left + bar_w + px(10), y, val_w, row_h,
                    f"{sc['aipmb_pct']:.0f}", BODY_SZ, ACCENT, True, PP_ALIGN.RIGHT)
        add_textbox(slide, bg_left + bar_w + val_w + px(5), y, val_w, row_h,
                    f"{sc['trad_pct']:.0f}", Pt(12), TEXT_GRAY, False, PP_ALIGN.RIGHT)


# ============================================================
# 流程图 — 居中
# ============================================================
def add_flow_steps(slide, steps, top_emu=None):
    """水平流程图，居中"""
    if top_emu is None:
        top_emu = BODY_TOP_EMU

    n = len(steps)
    step_w = px(200)
    arrow_w = px(30)
    step_h = px(66)
    total_w = n * step_w + (n - 1) * arrow_w
    start_left = center_x(total_w)

    for i, step in enumerate(steps):
        left = start_left + i * (step_w + arrow_w)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top_emu, step_w, step_h)

        color_map = {"active": ACCENT, "done": GREEN, "pending": GRAY_200}
        bg_c = color_map.get(step.get("color", "pending"), GRAY_200)
        text_c = WHITE if step.get("color") != "pending" else TEXT_GRAY

        set_shape_fill(shape, bg_c)
        shape.line.fill.background()

        tb = slide.shapes.add_textbox(left, top_emu, step_w, step_h)
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step.get("text", "")
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = text_c
        run.font.name = FONT_FAMILY

        if step.get("sub"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = step["sub"]
            run2.font.size = Pt(10)
            run2.font.color.rgb = text_c
            run2.font.name = FONT_FAMILY

        if i < n - 1:
            a_left = left + step_w
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                a_left + px(4), top_emu + (step_h - px(20)) // 2,
                px(22), px(20))
            set_shape_fill(arrow, GRAY_200)
            arrow.line.fill.background()


# ============================================================
# 金字塔 — 居中
# ============================================================
def add_pyramid(slide, levels, top_emu=None):
    """三层金字塔图，居中"""
    if top_emu is None:
        top_emu = BODY_TOP_EMU

    max_w = px(800)
    level_h = px(66)
    gap = px(6)

    for i, lv in enumerate(levels):
        w = int(max_w * lv.get("width_pct", 0.3))
        left = center_x(w)
        y = top_emu + i * (level_h + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, y, w, level_h)
        set_shape_fill(shape, lv.get("color", ACCENT))
        shape.line.fill.background()

        tb = slide.shapes.add_textbox(left, y, w, level_h)
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = lv.get("text", "")
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT_FAMILY

        if lv.get("sub"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = lv["sub"]
            run2.font.size = Pt(10)
            run2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            run2.font.name = FONT_FAMILY


# ============================================================
# 标签条/徽章
# ============================================================
def add_badge(slide, left_emu, top_emu, text, color=ACCENT):
    """单个彩色徽章"""
    w = px(len(text) * 10 + 30)
    h = px(26)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_emu, top_emu, w, h)
    set_shape_fill(s, ACCENT_LIGHT)
    s.line.fill.background()
    add_textbox(slide, left_emu, top_emu, w, h,
                text, Pt(10), color, False, PP_ALIGN.CENTER)


def add_badge_row(slide, badges, top_emu):
    """一行徽章，居中"""
    gap = px(8)
    badges_detail = []
    total_w = 0
    for text, color in badges:
        w = px(len(text) * 10 + 30)
        badges_detail.append((text, color, w))
        total_w += w
    total_w += (len(badges) - 1) * gap
    left = center_x(total_w)
    current_x = left
    for text, color, w in badges_detail:
        add_badge(slide, current_x, top_emu, text, color)
        current_x += w + gap


# ============================================================
# 两栏布局 — 居中
# ============================================================
def two_col_info(gap_px=30):
    """返回居中两栏坐标"""
    gap = px(gap_px)
    half_w = (CONTENT_W_EMU - gap) // 2
    left_x = center_x(CONTENT_W_EMU)
    right_x = left_x + half_w + gap
    return {"left_x": left_x, "right_x": right_x, "width": half_w}
