"""主生成逻辑 — 逐页构建 32 张幻灯片"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN

from .theme import *
from .components import *


def build_all_slides(prs, blank_layout):
    """构建全部 32 张幻灯片"""
    # ================================================================
    # Part 1: 项目概述 (P1-P4 + Divider)
    # ================================================================
    build_p1(prs, blank_layout)
    build_p2(prs, blank_layout)
    build_p3(prs, blank_layout)
    build_p4(prs, blank_layout)
    build_divider(prs, blank_layout, "02", "AI 原生架构设计")

    # ================================================================
    # Part 2: AI 原生架构设计 (P5-P12 + Divider)
    # ================================================================
    build_p5(prs, blank_layout)
    build_p6(prs, blank_layout)
    build_p7(prs, blank_layout)
    build_p8(prs, blank_layout)
    build_p9(prs, blank_layout)
    build_p10(prs, blank_layout)
    build_p11(prs, blank_layout)
    build_p12(prs, blank_layout)
    build_divider(prs, blank_layout, "03", "记忆与上下文系统")

    # ================================================================
    # Part 3: 记忆系统 (P13-P17 + Divider)
    # ================================================================
    build_p13(prs, blank_layout)
    build_p14(prs, blank_layout)
    build_p15(prs, blank_layout)
    build_p16(prs, blank_layout)
    build_p17(prs, blank_layout)
    build_divider(prs, blank_layout, "04", "前端体验设计")

    # ================================================================
    # Part 4: 前端体验 (P18-P22 + Divider)
    # ================================================================
    build_p18(prs, blank_layout)
    build_p19(prs, blank_layout)
    build_p20(prs, blank_layout)
    build_p21(prs, blank_layout)
    build_p22(prs, blank_layout)
    build_divider(prs, blank_layout, "05", "运营可观测体系")

    # ================================================================
    # Part 5: 运营 (P23-P27 + Divider)
    # ================================================================
    build_p23(prs, blank_layout)
    build_p24(prs, blank_layout)
    build_p25(prs, blank_layout)
    build_p26(prs, blank_layout)
    build_p27(prs, blank_layout)
    build_divider(prs, blank_layout, "06", "技术亮点与对比")

    # ================================================================
    # Part 6: 技术亮点 (P28-P30 + Divider)
    # ================================================================
    build_p28(prs, blank_layout)
    build_p29(prs, blank_layout)
    build_p30(prs, blank_layout)
    build_divider(prs, blank_layout, "07", "总结")

    # ================================================================
    # Part 7: 总结 (P31-P32)
    # ================================================================
    build_p31(prs, blank_layout)
    build_p32(prs, blank_layout)


# ================================================================
# Part 1: 项目概述
# ================================================================

def build_p1(prs, layout):
    """P1 封面"""
    slide = prs.slides.add_slide(layout)
    add_cover_slide(slide,
                    "AI 手机银行 11.0",
                    "AI 原生移动银行架构设计与技术方案",
                    "AIPmb 项目组 | 2026 年 6 月",
                    "面向 AI First 的银行服务架构实践")

def build_p2(prs, layout):
    """P2 为什么需要 AI 原生银行"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "为什么需要 AI 原生银行？")

    cols = two_col_info()
    y = BODY_TOP_EMU

    # 左栏：痛点
    add_textbox(slide, cols["left_x"], y, cols["width"], px(24),
                "传统银行 App 的痛点", H4_SZ, TEXT_GRAY, True)
    pain_points = [
        "功能堆砌，入口层级深，用户迷失在菜单中",
        "每个功能独立页面，上下文无法贯通",
        "搜索框只能匹配关键词，无法理解意图",
        "规则推荐，千人一面",
        "新增功能 = 新页面 + 新入口，边际成本递增",
    ]
    ptb = add_textbox(slide, cols["left_x"], y + px(26), cols["width"], px(180))
    for item in pain_points:
        add_paragraph(ptb.text_frame, "▸ " + item, SMALL_SZ, TEXT_DARK)

    # 右栏：优势
    add_textbox(slide, cols["right_x"], y, cols["width"], px(24),
                "AI 原生银行的优势", H4_SZ, ACCENT, True)
    advantages = [
        "自然语言统一入口：一句话直达任意服务",
        "上下文贯通：对话即会话，历史即记忆",
        "意图理解：LLM 深度解析用户需求",
        "千人千面：基于记忆的个人化服务",
        "功能即 Skill：新增能力只需注册 Skill，零前端改动",
    ]
    atb = add_textbox(slide, cols["right_x"], y + px(26), cols["width"], px(180))
    for item in advantages:
        add_paragraph(atb.text_frame, "▸ " + item, SMALL_SZ, TEXT_DARK)

    # 底部引用
    quote_y = y + px(250)
    quote = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_EMU, quote_y,
        SLIDE_W_EMU - 2 * MARGIN_EMU, px(40))
    set_shape_fill(quote, ACCENT_LIGHT)
    set_shape_border(quote, ACCENT, 0.5)
    add_textbox(slide, MARGIN_EMU + px(12), quote_y + px(8),
                SLIDE_W_EMU - 2 * MARGIN_EMU - px(24), px(24),
                "\"AI 原生银行不是'App 里加个聊天框'，而是以 LLM 为核心重新设计服务架构\"",
                SMALL_SZ, ACCENT, False, PP_ALIGN.LEFT)


def build_p3(prs, layout):
    """P3 项目全景"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "AIPmb 项目全景")

    layers = [
        {"label": "前端层", "bg_color": ARCH_BLUE_BG, "boxes": [
            {"text": "猎豹移动平台", "sub": "iOS / Android", "color": ACCENT_LIGHT, "text_color": PURPLE},
            {"text": "猎豹移动平台 Web", "sub": "运营后台", "color": ACCENT_LIGHT, "text_color": PURPLE},
        ]},
        {"label": "API 网关层", "bg_color": ARCH_CORE_BG, "boxes": [
            {"text": "FastAPI", "sub": "REST + WebSocket", "color": ACCENT_LIGHT, "text_color": ACCENT},
        ]},
        {"label": "AI 核心层", "bg_color": ARCH_DOMAIN_BG, "boxes": [
            {"text": "Agent 系统", "sub": "小易", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
            {"text": "Skill 编排", "sub": "12 个 Skill", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
            {"text": "LLM 适配器", "sub": "Kimi K2.5", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
            {"text": "记忆系统", "sub": "三层记忆", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
        ]},
        {"label": "数据服务层", "bg_color": ARCH_DATA_BG, "boxes": [
            {"text": "账户服务", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
            {"text": "交易服务", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
            {"text": "消费服务", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
            {"text": "产品服务", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
        ]},
        {"label": "数据层", "bg_color": ARCH_INFRA_BG, "boxes": [
            {"text": "Excel 数据源", "sub": "coredatas/", "color": RGBColor(0xF3, 0xE8, 0xFF), "text_color": PURPLE},
            {"text": "JSON 记忆", "sub": "memories/", "color": RGBColor(0xF3, 0xE8, 0xFF), "text_color": PURPLE},
        ]},
    ]
    add_arch_diagram(slide, layers)


def build_p4(prs, layout):
    """P4 核心创新点"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "核心创新点")

    cards = [
        {"icon": "🎯", "title": "LLM as Controller",
         "desc": "通用助手\"小易\"作为唯一智能入口，LLM 通过 Function Calling 自主决策调用 Skill"},
        {"icon": "🧩", "title": "双层 Skill 编排",
         "desc": "领域层 Skill（理财/消费/画像等7个）+ 数据层 Skill（账户/交易/产品等5个），能力复用灵活组合"},
        {"icon": "⚡", "title": "流式 Thinking 可视化",
         "desc": "四阶段思考过程（意图识别→Skill编排→数据收集→分析生成）实时推送前端"},
        {"icon": "🧠", "title": "推理链透传",
         "desc": "Kimi reasoning_content 字段透传，前端展示 AI 内部推理过程，增强信任感和透明度"},
        {"icon": "💾", "title": "三层记忆架构",
         "desc": "会话级（20轮上下文）+ 用户级（JSON持久化）+ 压缩级（LLM自动摘要）"},
        {"icon": "📸", "title": "多模态预留",
         "desc": "vision() 接口 + image_picker 已集成，支持图片输入，模型不支持时自动降级"},
    ]
    add_card_grid(slide, cards, cols=2)


# ================================================================
# Part 2: AI 原生架构设计
# ================================================================

def build_p5(prs, layout):
    """P5 传统 vs AI 原生"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "传统架构 vs AI 原生架构")

    headers = ["维度", "传统银行 App", "AI 原生银行（AIPmb）"]
    cw = px(200)
    cw2 = px(540)
    rows = [
        {"cells": ["用户入口", "多级菜单 + 搜索框", "自然语言对话，一句话直达任意服务"], "highlight_col": 2},
        {"cells": ["服务发现", "用户主动寻找功能", "AI 主动理解意图，推荐服务"], "highlight_col": 2},
        {"cells": ["数据查询", "固定报表页面", "LLM 动态生成分析，任意维度组合"], "highlight_col": 2},
        {"cells": ["上下文贯通", "页面跳转即丢失上下文", "对话即会话，历史记忆自动关联"], "highlight_col": 2},
        {"cells": ["个性化", "规则引擎推荐", "基于记忆的理解式推荐，千人千面"], "highlight_col": 2},
        {"cells": ["扩展性", "新增功能=新增页面+入口", "新增功能=新增 Skill，零前端改动"], "highlight_col": 2},
        {"cells": ["开发模式", "前端页面驱动", "LLM 编排驱动"], "highlight_col": 2},
    ]
    add_table(slide, headers, rows, [cw, cw2, cw2])


def build_p6(prs, layout):
    """P6 LLM as Controller"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "LLM as Controller：统一智能入口")

    cols = two_col_info()
    left_x = cols["left_x"]
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, left_x, y, cw, px(300))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "\"小易\"作为单一智能体入口", Pt(13), ACCENT, True)
    add_paragraph(tf, "LLM 通过 Function Calling 自主决策调用 Skill", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "AgentRegistry 评分阈值 0.3，低于阈值自动路由到小易", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "预留 @智能体名 显式指定扩展能力", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "设计理念：银行场景意图相对明确，单一 Agent + Skill 编排比多 Agent 协作更简洁高效", Pt(10), TEXT_GRAY)

    # 右栏代码
    right_x = cols["right_x"]
    code = """# AgentRegistry 路由
class AgentRegistry:
    THRESHOLD = 0.3

    async def route(self, query, agents, llm):
        scores = await llm.score(query, agents)
        best = max(scores, key=scores.get)
        if scores[best] >= self.THRESHOLD:
            return best
        return "xiaoyi"  # 兜底"""
    add_code_block(slide, code, right_x, y, cw, px(180), "Python")


def build_p7(prs, layout):
    """P7 双层 Skill 编排"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "双层 Skill 编排架构")

    layers = [
        {"label": "用户输入", "bg_color": ACCENT_LIGHT, "boxes": [
            {"text": "帮我制定理财方案", "color": ACCENT, "text_color": WHITE},
        ]},
        {"label": "领域层 Skill（7个）", "bg_color": ARCH_DOMAIN_BG, "boxes": [
            {"text": "理财规划", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
            {"text": "消费分析", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
            {"text": "收入预测", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
            {"text": "用户画像", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
            {"text": "生活推荐", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
            {"text": "社区分析", "color": RGBColor(0xDB, 0xEA, 0xFE), "text_color": BLUE},
        ]},
        {"label": "数据层 Skill（5个）", "bg_color": ARCH_DATA_BG, "boxes": [
            {"text": "账户查询", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
            {"text": "交易查询", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
            {"text": "消费统计", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
            {"text": "产品查询", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
            {"text": "标签计算", "color": RGBColor(0xD1, 0xFA, 0xE5), "text_color": GREEN},
        ]},
    ]
    add_arch_diagram(slide, layers)


def build_p8(prs, layout):
    """P8 Skill 编排实战"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "Skill 编排实战：理财规划")

    y = BODY_TOP_EMU
    steps = [
        {"text": "1. 意图识别", "sub": "理财规划", "color": "done"},
        {"text": "2. Skill编排", "sub": "6个数据Skill", "color": "active"},
        {"text": "3. 数据收集", "sub": "并行调用", "color": "pending"},
        {"text": "4. 分析生成", "sub": "LLM汇总", "color": "pending"},
    ]
    add_flow_steps(slide, steps, y)

    y2 = y + px(80)
    cols = two_col_info()
    cw = cols["width"]

    # 左：说明
    tb = add_textbox(slide, cols["left_x"], y2, cw, px(200))
    tf = tb.text_frame
    tf.paragraphs[0].text = "用户说：\"帮我制定一个理财方案\""
    tf.paragraphs[0].font.size = SMALL_SZ
    tf.paragraphs[0].font.color.rgb = TEXT_GRAY
    add_paragraph(tf, "6 个数据 Skill 并行调用：", SMALL_SZ, TEXT_DARK, True)
    skills = ["collect_account_data", "collect_consumption_data", "collect_transaction_data",
              "collect_product_data", "compute_user_tags", "calculate_survival"]
    for s in skills:
        add_paragraph(tf, "  ▸ " + s, CAPTION_SZ, TEXT_GRAY)

    # 右：JSON 结果
    json_code = """{
  "total_balance": 63750.31,
  "monthly_income": 12000,
  "monthly_expense": 3500,
  "risk_profile": "稳健型",
  "survival_months": 18,
  "user_tags": ["工薪族", "有储蓄"]
}"""
    add_code_block(slide, json_code, cols["right_x"], y2, cw, px(200), "JSON")


def build_p9(prs, layout):
    """P9 流式 Thinking 可视化"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "流式 Thinking 可视化")

    y = BODY_TOP_EMU
    steps = [
        {"text": "🎯 意图识别", "sub": "理解用户需求", "color": "done"},
        {"text": "🧩 Skill 编排", "sub": "选择合适Skill", "color": "active"},
        {"text": "📊 数据收集", "sub": "并行查询数据", "color": "pending"},
        {"text": "✨ 分析生成", "sub": "LLM生成回复", "color": "pending"},
    ]
    add_flow_steps(slide, steps, y)

    y2 = y + px(90)
    tb = add_textbox(slide, MARGIN_EMU, y2, SLIDE_W_EMU - 2 * MARGIN_EMU, px(120))
    tf = tb.text_frame
    add_paragraph(tf, "事件类型定义", SMALL_SZ, TEXT_DARK, True)
    events = [
        "thinking_step — 思考阶段进度推送",
        "reasoning_chunk — Kimi reasoning_content 透传",
        "ai_chunk — LLM 流式输出文本",
        "ai_done — 回复完成信号",
        "error — 异常事件",
    ]
    for ev in events:
        add_paragraph(tf, "  ▸ " + ev, CAPTION_SZ, TEXT_DARK)

    # Badge 行
    badges = [("thinking_step", ACCENT), ("reasoning_chunk", BLUE),
              ("ai_chunk", GREEN), ("ai_done", PURPLE), ("error", ORANGE)]
    add_badge_row(slide, badges, y2 + px(140))


def build_p10(prs, layout):
    """P10 推理链透传"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "推理链透传：看见 AI 的思考")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(200))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "推理链四大价值", Pt(13), ACCENT, True)
    add_paragraph(tf, "▸ 增强信任：用户看到 AI 推理过程", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "▸ 可解释性：理解 AI 为什么这样决策", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "▸ 调试优化：开发者可定位问题", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "▸ 教育价值：用户从推理中学习理财知识", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "Kimi K2.5 支持 reasoning_content 字段，包含模型内部思维链，AIPmb 完整透传到前端", Pt(10), TEXT_GRAY)

    # 右：终端 mockup
    term = """💭 推理过程
用户询问理财方案，我需要：
1. 先了解用户的资产状况
2. 分析消费习惯和风险承受能力
3. 从产品库筛选匹配的产品
4. 生成个性化的配置建议
→ 调用 financial_planning Skill
📋 基于您的资产状况（总资产 ¥63,750）
建议采用"稳健型"配置方案..."""
    add_code_block(slide, term, cols["right_x"], y, cw, px(230), "Reasoning")


def build_p11(prs, layout):
    """P11 代码示例"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "代码示例：通用助手的流式分析")

    code = """# pmb/agents/general_assistant.py
async def analyze_stream(self, context, event_queue):
    conv = context_manager.get_or_create(context.session_id)
    messages = conv.get_messages()

    # 阶段 1: 意图识别
    await event_queue.put(ThinkingEvent(
        step="意图识别", status="active"))
    intent = await self.llm.identify_intent(messages)

    # 阶段 2: Skill 编排
    await event_queue.put(ThinkingEvent(
        step="Skill 编排", status="active"))
    skills = self.skill_orchestrator.select(intent)

    # 阶段 3: 数据收集（并行）
    await event_queue.put(ThinkingEvent(
        step="数据收集", status="active"))
    data = await asyncio.gather(*[
        skill.execute(context) for skill in skills])

    # 阶段 4: 分析生成（流式输出）
    await event_queue.put(ThinkingEvent(
        step="分析生成", status="active"))
    async for chunk in self.llm.chat_stream(
            self._build_prompt(intent, data)):
        await event_queue.put(AiChunkEvent(chunk))"""
    add_code_block(slide, code, MARGIN_EMU, BODY_TOP_EMU,
                   SLIDE_W_EMU - 2 * MARGIN_EMU,
                   SLIDE_H_EMU - BODY_TOP_EMU - MARGIN_EMU, "Python")


def build_p12(prs, layout):
    """P12 架构设计原则"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "架构设计原则")

    cards = [
        {"icon": "1️⃣", "title": "单一智能体入口",
         "desc": "小易作为唯一入口，通过 Function Calling 智能路由到 12 个 Skill，避免多 Agent 复杂度"},
        {"icon": "2️⃣", "title": "Skill 无状态纯函数",
         "desc": "每个 Skill 独立/无状态/可测试，像微服务一样注册编排，支持灵活组合"},
        {"icon": "3️⃣", "title": "事件驱动流式响应",
         "desc": "asyncio.Queue + 四种事件类型，前端实时展示 AI 思考和执行全过程"},
        {"icon": "4️⃣", "title": "LLM 只做决策不做数据",
         "desc": "LLM 负责意图识别和结果生成，数据采集和计算由 Skill/Tool 完成，分离关注点"},
    ]
    add_card_grid(slide, cards, cols=2)


# ================================================================
# Part 3: 记忆与上下文
# ================================================================

def build_p13(prs, layout):
    """P13 为什么记忆是关键"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "为什么记忆是 AI 银行的关键？")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    # 左：无记忆
    box_h = px(200)
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cols["left_x"], y, cw, box_h)
    set_shape_fill(b1, GRAY_100)
    set_shape_border(b1, GRAY_200, 0.5)
    add_textbox(slide, cols["left_x"] + px(12), y + px(6), cw - px(24), px(18),
                "没有记忆的场景", SMALL_SZ, TEXT_GRAY, True)
    t1 = add_textbox(slide, cols["left_x"] + px(12), y + px(30), cw - px(24), box_h - px(36))
    add_paragraph(t1.text_frame, "用户（昨天）：\"我想了解一下理财产品\"", CAPTION_SZ, TEXT_DARK)
    add_paragraph(t1.text_frame, "小易：\"好的，为您推荐以下...\"", CAPTION_SZ, BLUE)
    add_paragraph(t1.text_frame, "", CAPTION_SZ, TEXT_DARK)
    add_paragraph(t1.text_frame, "用户（今天）：\"再详细说说那个方案\"", CAPTION_SZ, TEXT_DARK)
    add_paragraph(t1.text_frame, "小易：\"请问您指的是哪个方案？\" ❌", CAPTION_SZ, ORANGE)

    # 右：有记忆
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cols["right_x"], y, cw, box_h)
    set_shape_fill(b2, ACCENT_LIGHT)
    set_shape_border(b2, ACCENT, 0.5)
    add_textbox(slide, cols["right_x"] + px(12), y + px(6), cw - px(24), px(18),
                "有记忆的场景", SMALL_SZ, ACCENT, True)
    t2 = add_textbox(slide, cols["right_x"] + px(12), y + px(30), cw - px(24), box_h - px(36))
    add_paragraph(t2.text_frame, "用户（昨天）：\"我想了解一下理财产品\"", CAPTION_SZ, TEXT_DARK)
    add_paragraph(t2.text_frame, "小易：\"好的，为您推荐以下...\"", CAPTION_SZ, BLUE)
    add_paragraph(t2.text_frame, "", CAPTION_SZ, TEXT_DARK)
    add_paragraph(t2.text_frame, "用户（今天）：\"再详细说说那个方案\"", CAPTION_SZ, TEXT_DARK)
    add_paragraph(t2.text_frame, "小易：\"好的，您昨天了解的稳健型基金方案，最新净值...\" ✅", CAPTION_SZ, GREEN)

    # 底部引用
    qy = y + box_h + px(12)
    add_textbox(slide, MARGIN_EMU, qy, SLIDE_W_EMU - 2 * MARGIN_EMU, px(24),
                "\"记忆是 AI 从'工具'升级为'伙伴'的关键\"", SMALL_SZ, ACCENT, False, PP_ALIGN.CENTER)


def build_p14(prs, layout):
    """P14 三层记忆架构"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "三层记忆架构")

    levels = [
        {"text": "压缩级记忆", "sub": "LLM 自动摘要 ≤150字 | 30天自动清理",
         "width_pct": 0.35, "color": PURPLE},
        {"text": "用户级记忆", "sub": "跨会话长期存储 | JSON 文件持久化 | 用户粒度并发锁",
         "width_pct": 0.62, "color": BLUE},
        {"text": "会话级记忆", "sub": "20 轮上下文窗口 | 内存实时对话 | 自动 Trim 保留最近 N 轮",
         "width_pct": 0.92, "color": ACCENT},
    ]
    add_pyramid(slide, levels)

    # Badge 行
    y3 = BODY_TOP_EMU + 3 * px(56)
    badges = [("时效性 ↑", ACCENT), ("容量 ↑", BLUE), ("持久性 ↑", GREEN)]
    add_badge_row(slide, badges, y3)


def build_p15(prs, layout):
    """P15 会话级记忆"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "会话级记忆：20 轮上下文窗口")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(200))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "ContextManager 设计", Pt(13), ACCENT, True)
    add_paragraph(tf, "按 session_id 管理独立会话", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "自动 Trim 保留最近 20 轮消息", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "System Prompt 动态注入用户记忆摘要", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "会话期间全内存操作，零 I/O 延迟", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "设计考量：20 轮是成本与效果的平衡点。超过 20 轮后 LLM 注意力分散，Token 成本急剧上升", CAPTION_SZ, TEXT_GRAY)

    code = """class ConversationContext:
    MAX_MESSAGES = 20

    def add_message(self, role, content):
        self.messages.append({
            "role": role,
            "content": content,
            "time": time.time()
        })
        # 自动 Trim
        if len(self.messages) > self.MAX_MESSAGES:
            self.messages = self.messages[
                -self.MAX_MESSAGES:]"""
    add_code_block(slide, code, cols["right_x"], y, cw, px(200), "Python")


def build_p16(prs, layout):
    """P16 用户级记忆"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "用户级记忆：跨会话持久化")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(250))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "MemoryService 生命周期", Pt(13), ACCENT, True)
    add_paragraph(tf, "start_session: 新建会话并加载历史摘要", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "end_session: 压缩对话并保存到 JSON 文件", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "JSON 路径: output/memories/{user}.json", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "asyncio.Lock 按 user_name 加锁，避免并发写入冲突", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "注入流程:", SMALL_SZ, TEXT_DARK, True)
    add_paragraph(tf, "连接 WebSocket → start_session → 加载摘要 → 注入 System Prompt → LLM 生成回复",
                  CAPTION_SZ, TEXT_GRAY)

    json_code = """{
  "user_name": "彭楫洲",
  "entries": [
    {"session_id": "sess_1",
     "summary": "咨询稳健型基金",
     "time": "2026-05-20",
     "messages_count": 12},
    {"session_id": "sess_2",
     "summary": "分析5月消费情况",
     "time": "2026-05-21",
     "messages_count": 8}
  ],
  "last_updated": "2026-05-21"
}"""
    add_code_block(slide, json_code, cols["right_x"], y, cw, px(250), "JSON")


def build_p17(prs, layout):
    """P17 压缩级记忆"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "压缩级记忆：LLM 智能摘要")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(200))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "自动压缩策略", Pt(13), ACCENT, True)
    add_paragraph(tf, "会话结束后超过 1 天触发压缩", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "LLM 压缩为 ≤150 字摘要", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "超过 30 天的旧记忆自动清理", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "LLM 不可用时保留原始数据，不阻塞", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "压缩前：12 条消息，约 3000 字", Pt(10), TEXT_GRAY)
    add_paragraph(tf, "压缩后：\"用户咨询理财产品，偏好稳健型，总资产约 6.4 万\" (25字)", Pt(10), ACCENT)

    code = """async def compress_if_needed(self, user, session):
    if not self._should_compress(session):
        return
    summary = await self.llm.compress(
        session.get_messages(),
        max_length=150
    )
    await self.storage.save_summary(
        user_name=user,
        session_id=session.id,
        summary=summary
    )"""
    add_code_block(slide, code, cols["right_x"], y, cw, px(200), "Python")


# ================================================================
# Part 4: 前端体验设计
# ================================================================

def build_p18(prs, layout):
    """P18 猎豹移动平台跨端方案"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "猎豹移动平台 跨端方案")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(250))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "技术选型", Pt(13), ACCENT, True)
    add_paragraph(tf, "猎豹移动平台：一套代码覆盖 iOS/Android/Web", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "Riverpod 状态管理：编译时安全", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "WebSocket 实时通信", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "项目结构", SMALL_SZ, TEXT_DARK, True)
    add_paragraph(tf, "aipmb_app/ — 客户端", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf, "aipmb_manage/ — 运营后台", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf, "共享 models/providers/services", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "核心优势：业务逻辑 100% 复用，UI 层针对不同平台微调", Pt(10), ACCENT)

    # 右：三端 mockup
    rx = cols["right_x"]
    for i, label in enumerate(["iOS", "Android", "Web"]):
        mw = px(160)
        mh = px(300)
        ml = rx + i * (mw + px(12))
        phone = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, ml, y, mw, mh)
        set_shape_fill(phone, WHITE)
        set_shape_border(phone, DARK, 1.5)
        add_textbox(slide, ml, y + mh // 2 - px(10), mw, px(20),
                    label, SMALL_SZ, TEXT_GRAY, False, PP_ALIGN.CENTER)


def build_p19(prs, layout):
    """P19 此刻页面"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "\"此刻\"页面：AI 驱动的智能首页")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(250))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "设计理念", Pt(13), ACCENT, True)
    add_paragraph(tf, "去菜单化：无九宫格，AI 入口前置", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "智能推荐卡片：基于用户画像动态推荐", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "实时数据展示：资产总览一目了然", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "页面结构", SMALL_SZ, TEXT_DARK, True)
    add_paragraph(tf, "顶部 — 资产总览", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "中部 — AI 对话入口", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "底部 — 推荐卡片", SMALL_SZ, TEXT_DARK)

    # 右：手机 mockup
    rx = cols["right_x"]
    mw = px(200)
    mh = px(360)
    phone = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, rx, y, mw, mh)
    set_shape_fill(phone, WHITE)
    set_shape_border(phone, DARK, 1.5)

    # Mockup 内文字
    atb = add_textbox(slide, rx + px(14), y + px(24), mw - px(28), mh - px(48))
    tf2 = atb.text_frame
    add_paragraph(tf2, "总资产", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf2, "¥63,750.31", Pt(16), TEXT_DARK, True)
    add_paragraph(tf2, "", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf2, "你好，彭楫洲", SMALL_SZ, ACCENT, True)
    add_paragraph(tf2, "有什么可以帮你的？", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf2, "", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf2, "▸ 理财方案", CAPTION_SZ, TEXT_DARK)
    add_paragraph(tf2, "▸ 消费分析", CAPTION_SZ, TEXT_DARK)
    add_paragraph(tf2, "▸ 贷款咨询", CAPTION_SZ, TEXT_DARK)
    add_paragraph(tf2, "", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf2, "今日支出 ¥328.50", CAPTION_SZ, TEXT_GRAY)


def build_p20(prs, layout):
    """P20 全屏对话"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "全屏对话：与小易的自然交互")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(250))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "交互设计", Pt(13), ACCENT, True)
    add_paragraph(tf, "全屏沉浸式对话界面", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "消息气泡（用户右对齐/AI 左对齐）", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "Thinking 四阶段进度动画", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "推理链默认折叠，点击可展开", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "推荐问题卡片，一键发送", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "消息类型", SMALL_SZ, TEXT_DARK, True)
    badges = [("用户消息", ACCENT), ("AI 回复", BLUE), ("Thinking", GREEN),
              ("推理链", PURPLE), ("Tool 结果", ORANGE)]
    add_badge_row(slide, badges, y + px(230))

    # 右：对话 mockup
    rx = cols["right_x"]
    mw = px(200)
    mh = px(380)
    phone = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, rx, y, mw, mh)
    set_shape_fill(phone, WHITE)
    set_shape_border(phone, DARK, 1.5)

    atb = add_textbox(slide, rx + px(12), y + px(16), mw - px(24), mh - px(32))
    tf2 = atb.text_frame
    add_paragraph(tf2, "小易", CAPTION_SZ, TEXT_GRAY, True)
    add_paragraph(tf2, "", Pt(8), TEXT_GRAY)
    add_paragraph(tf2, "帮我看看怎么理财", CAPTION_SZ, ACCENT, False)
    add_paragraph(tf2, "───────────────", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf2, "意图识别 → Skill编排 → 数据收集 → 分析生成", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf2, "正在分析您的资产...", CAPTION_SZ, TEXT_DARK)
    add_paragraph(tf2, "", Pt(8), TEXT_GRAY)
    add_paragraph(tf2, "推理过程", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf2, "用户总资产 6.4 万，消费习惯稳健，建议配置40%货币基金+30%债券+30%混合基金", CAPTION_SZ, TEXT_DARK)


def build_p21(prs, layout):
    """P21 流式 State 管理"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "流式响应的 State 管理")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(200))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "Riverpod 架构", Pt(13), ACCENT, True)
    add_paragraph(tf, "chatMessagesProvider — 消息列表管理", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "chatStatusProvider — idle/loading/streaming/error", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "thinkingStepsProvider — 四阶段进度追踪", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "WebSocketService 全局单例，避免重复连接", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "关键设计：StreamController 懒创建且永不关闭，页面切换时 Notifier 重新监听", CAPTION_SZ, TEXT_GRAY)

    code = """class ChatMessagesNotifier
    extends StateNotifier<List<ChatMessage>> {

  void _onEvent(dynamic event) {
    if (event is ThinkingStepEvent) {
      state = [...state, event.toThinking()];
    } else if (event is AiChunkEvent) {
      // 流式追加当前消息
      state = state.last is AiMessage
        ? [...state.sublist(0, len-1),
           state.last.copyWith(
             content: state.last.content
               + event.chunk)]
        : [...state, AiMessage(event.chunk)];
    }
  }
}"""
    add_code_block(slide, code, cols["right_x"], y, cw, px(280), "Dart")


def build_p22(prs, layout):
    """P22 多模态预留"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "多模态预留：图片输入能力")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(250))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "架构设计", Pt(13), ACCENT, True)
    add_paragraph(tf, "vision() 接口：统一多模态调用", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "image_picker 已集成", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "Base64 编码传输", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "模型不支持时自动降级为文本描述", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "应用场景", SMALL_SZ, TEXT_DARK, True)
    add_paragraph(tf, "▸ 拍摄银行卡识别卡号", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf, "▸ 上传消费小票自动记账", CAPTION_SZ, TEXT_GRAY)
    add_paragraph(tf, "▸ 拍摄证件 OCR 提取信息", CAPTION_SZ, TEXT_GRAY)

    code = """class QwenLLM(BaseLLM):
    async def vision(self, image_data,
                     prompt, model="qwen-vl"):
        if not self.supports_vision(model):
            return await self.chat(
                f"[图片描述]{prompt}")
        return await self._call_vision_api(
            model, image_data, prompt)"""
    add_code_block(slide, code, cols["right_x"], y, cw, px(180), "Python")


# ================================================================
# Part 5: 运营可观测体系
# ================================================================

def build_p23(prs, layout):
    """P23 AI-Manage 全景"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "AI-Manage：运营后台全景")

    y = BODY_TOP_EMU
    layers = [
        {"label": "运营后台功能模块", "bg_color": ACCENT_LIGHT, "boxes": [
            {"text": "用户管理", "sub": "查看/搜索用户", "color": ACCENT_LIGHT, "text_color": ACCENT},
            {"text": "对话审计", "sub": "查询对话记录", "color": ACCENT_LIGHT, "text_color": BLUE},
            {"text": "业务标注", "sub": "AI 自动维度识别", "color": ACCENT_LIGHT, "text_color": GREEN},
            {"text": "模型配置", "sub": "切换 LLM 后端", "color": ACCENT_LIGHT, "text_color": PURPLE},
        ]},
    ]
    add_arch_diagram(slide, layers, y)

    cols = two_col_info()
    cw = cols["width"]
    y2 = y + px(80)

    # 左：技术实现
    t1 = add_textbox(slide, cols["left_x"], y2, cw, px(100))
    tf1 = t1.text_frame
    add_paragraph(tf1, "技术实现", SMALL_SZ, ACCENT, True)
    add_paragraph(tf1, "猎豹移动平台 Web + REST API", CAPTION_SZ, TEXT_DARK)
    add_paragraph(tf1, "实时内存会话 + 历史 JSON 文件聚合", CAPTION_SZ, TEXT_DARK)

    # 右：核心价值
    t2 = add_textbox(slide, cols["right_x"], y2, cw, px(100))
    tf2 = t2.text_frame
    add_paragraph(tf2, "核心价值", SMALL_SZ, ACCENT, True)
    add_paragraph(tf2, "服务质量监控 · 用户需求洞察 · 合规审计 · 运营决策支撑", CAPTION_SZ, TEXT_DARK)


def build_p24(prs, layout):
    """P24 对话审计"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "对话审计：看见每一次 AI 交互")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(200))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "四维查询", Pt(13), ACCENT, True)
    add_paragraph(tf, "时间维度：today / week / month", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "业务维度：理财/贷款/消费/账户/其他", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "关键词搜索：消息内容全文检索", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "用户维度：从用户详情一键查看全部对话", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "数据聚合：内存会话 + JSON 文件 → 关键词预标注 → 按时间归类 → 过滤分页 → AI 分析",
                  CAPTION_SZ, TEXT_GRAY)

    term = """会话: sess_abc123 | 用户: 彭楫洲 | 业务: 理财
用户: 帮我制定理财方案
Thinking: 意图识别 → Skill编排
AI: 根据您的资产状况(¥63,750)...
用户: 有具体的产品推荐吗？
AI: 建议配置：40%货币基金..."""
    add_code_block(slide, term, cols["right_x"], y, cw, px(230), "Audit")


def build_p25(prs, layout):
    """P25 业务标注"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "业务标注：AI 自动维度识别")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(200))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "自动标注策略", Pt(13), ACCENT, True)
    add_paragraph(tf, "理财/基金/存款 → 理财维度", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "贷款/房贷/利率 → 贷款维度", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "消费/支出/账单 → 消费维度", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "账户/余额/转账 → 账户维度", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "未匹配 → 其他维度", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "未来演进：接入 LLM 语义分类，准确率可提升至 95%+", CAPTION_SZ, ACCENT)

    # 右：标注效果
    atb = add_textbox(slide, cols["right_x"], y, cw, px(200))
    tf2 = atb.text_frame
    add_paragraph(tf2, "标注效果示例", Pt(13), ACCENT, True)
    examples = [("理财", BLUE), ("贷款", ORANGE), ("消费", GREEN), ("账户", PURPLE)]
    for label, color in examples:
        add_badge(slide, cols["right_x"] + px(10), y + px(40) + examples.index((label, color)) * px(30),
                  "  " + label + "  " + "相关问题", color)
        y_shift = 0


def build_p26(prs, layout):
    """P26 模型配置"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "模型配置：灵活切换 LLM 后端")

    cols = two_col_info()
    cw = cols["width"]
    y = BODY_TOP_EMU

    tb = add_textbox(slide, cols["left_x"], y, cw, px(200))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    add_paragraph(tf, "适配器模式", Pt(13), ACCENT, True)
    add_paragraph(tf, "BaseLLM 抽象统一接口 chat/chat_stream/vision", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "多模型支持 Kimi/Qwen/GPT 无缝切换", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "配置化：环境变量或管理后台切换模型", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "reasoning_content 透传 + Function Calling 格式转换", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "", SMALL_SZ, TEXT_DARK)
    add_paragraph(tf, "当前主模型：Kimi K2.5", SMALL_SZ, ACCENT, True)
    add_paragraph(tf, "支持推理链：是 | Function Calling：是", CAPTION_SZ, TEXT_DARK)

    code = """class BaseLLM(ABC):
    @abstractmethod
    async def chat(self, msgs, **kw):
        pass
    @abstractmethod
    async def chat_stream(self, msgs, **kw):
        pass

class KimiLLM(BaseLLM):
    model = "kimi-k2.5"
    def _build_request(self, msgs):
        return {"model": self.model,
                "messages": msgs,
                "stream": True}"""
    add_code_block(slide, code, cols["right_x"], y, cw, px(220), "Python")


def build_p27(prs, layout):
    """P27 价值闭环"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "运营数据的价值闭环")

    y = BODY_TOP_EMU
    steps = [
        {"text": "1. 数据采集", "sub": "实时+历史对话", "color": "done"},
        {"text": "2. 业务标注", "sub": "AI 自动维度识别", "color": "active"},
        {"text": "3. 统计分析", "sub": "高频问题/分布", "color": "pending"},
        {"text": "4. 产品优化", "sub": "反馈至 Skill", "color": "pending"},
    ]
    add_flow_steps(slide, steps, y)

    y2 = y + px(80)
    cards = [
        {"icon": "📊", "title": "对话量统计", "desc": "日/周/月对话趋势，活跃用户数"},
        {"icon": "📈", "title": "业务分布", "desc": "理财/贷款/消费等维度占比分析"},
        {"icon": "⏱️", "title": "响应时长", "desc": "平均响应时间，流式首字节时间"},
        {"icon": "⭐", "title": "满意度追踪", "desc": "用户反馈收集与趋势分析"},
    ]
    add_card_grid(slide, cards, cols=2, top_emu=y2)


# ================================================================
# Part 6: 技术亮点
# ================================================================

def build_p28(prs, layout):
    """P28 架构先进性对比"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "架构先进性对比")

    scores = [
        {"label": "意图理解", "aipmb_pct": 95, "trad_pct": 30},
        {"label": "上下文贯通", "aipmb_pct": 90, "trad_pct": 20},
        {"label": "个性化程度", "aipmb_pct": 85, "trad_pct": 40},
        {"label": "扩展性", "aipmb_pct": 92, "trad_pct": 35},
        {"label": "开发效率", "aipmb_pct": 88, "trad_pct": 45},
        {"label": "可观测性", "aipmb_pct": 85, "trad_pct": 50},
    ]
    add_score_bars(slide, scores, BODY_TOP_EMU + px(20))

    # 图例
    ly = BODY_TOP_EMU + 6 * px(42) + px(30)
    leg_w = px(240)
    leg_left = center_x(leg_w)
    leg1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, leg_left, ly, px(16), px(16))
    set_shape_fill(leg1, ACCENT)
    leg1.line.fill.background()
    add_textbox(slide, leg_left + px(20), ly, px(80), px(16),
                "AIPmb", CAPTION_SZ, ACCENT, True)
    leg2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, leg_left + px(120), ly, px(16), px(16))
    set_shape_fill(leg2, GRAY_200)
    leg2.line.fill.background()
    add_textbox(slide, leg_left + px(140), ly, px(120), px(16),
                "传统银行 App", CAPTION_SZ, TEXT_GRAY)


def build_p29(prs, layout):
    """P29 代码质量"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "代码质量与工程实践")

    cards = [
        {"icon": "🔌", "title": "抽象接口设计",
         "desc": "BaseLLM/BaseSkill/BaseAgent 统一抽象，新增模型只需实现接口"},
        {"icon": "🔒", "title": "并发安全",
         "desc": "用户粒度 asyncio.Lock，按 user_name 加锁避免并发写入冲突"},
        {"icon": "🧪", "title": "可测试性",
         "desc": "Skill 纯函数设计 + LLM 可 Mock，单元测试覆盖核心逻辑"},
        {"icon": "📦", "title": "模块化架构",
         "desc": "agents/skills/llm/core 清晰分层，每层独立可替换"},
        {"icon": "⚡", "title": "状态管理",
         "desc": "Riverpod 编译时安全，Provider 自动销毁，避免内存泄漏"},
        {"icon": "🛡️", "title": "降级策略",
         "desc": "LLM 不可用不阻塞服务，压缩失败保留原始数据，图片自动降级"},
    ]
    add_card_grid(slide, cards, cols=2)


def build_p30(prs, layout):
    """P30 未来演进"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "未来演进路线")

    y = BODY_TOP_EMU + px(20)
    # 三圆居中：圆宽80，间距380
    circle_w = px(80)
    total_w = 3 * circle_w + px(760)
    c_start = center_x(total_w)
    phases = [
        {"title": "近期", "sub": "多模态上线 | 语音对话 | 更多 Skill",
         "color": ACCENT, "left": c_start},
        {"title": "中期", "sub": "多 Agent 协作 | 智能体市场 | 推荐引擎",
         "color": BLUE, "left": c_start + circle_w + px(380)},
        {"title": "远期", "sub": "自主决策 Agent | 跨银行生态 | 情感计算",
         "color": PURPLE, "left": c_start + 2 * (circle_w + px(380))},
    ]

    for ph in phases:
        # 圆形节点
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            ph["left"], y, px(80), px(80))
        set_shape_fill(node, ph["color"])
        node.line.fill.background()
        add_textbox(slide, ph["left"], y + px(24), px(80), px(32),
                    ph["title"], Pt(14), WHITE, True, PP_ALIGN.CENTER)
        # 副标题在下方
        add_textbox(slide, ph["left"] - px(30), y + px(90), px(140), px(40),
                    ph["sub"], CAPTION_SZ, TEXT_GRAY, False, PP_ALIGN.CENTER)

    # 连接线
    line_y = y + px(40)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        c_start + px(40), line_y - px(2), total_w - px(80), px(4))
    set_shape_fill(line, GRAY_200)
    line.line.fill.background()

    # 愿景
    vy = y + px(160)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_EMU, vy, SLIDE_W_EMU - 2 * MARGIN_EMU, px(56))
    set_shape_fill(box, ACCENT_LIGHT)
    set_shape_border(box, ACCENT, 1)
    add_textbox(slide, MARGIN_EMU + px(16), vy + px(8),
                SLIDE_W_EMU - 2 * MARGIN_EMU - px(32), px(40),
                "愿景：从\"AI 辅助银行服务\"演进为\"AI 原生银行服务\"",
                Pt(14), ACCENT, True, PP_ALIGN.CENTER)


# ================================================================
# Part 7: 总结
# ================================================================

def build_p31(prs, layout):
    """P31 核心结论"""
    slide = prs.slides.add_slide(layout)
    add_header(slide, "核心结论")

    y = BODY_TOP_EMU + px(20)
    conclusions = [
        ("01", ACCENT, "LLM as Controller 是 AI 原生银行的核心",
         "以 LLM 为统一入口，通过 Skill 编排完成复杂任务，取代传统菜单驱动"),
        ("02", BLUE, "记忆系统决定 AI 的服务上限",
         "三层记忆让 AI 从\"陌生人\"进化为\"私人顾问\""),
        ("03", GREEN, "流式可视化增强用户信任",
         "Thinking 过程实时展示 + 推理链透传，让 AI 决策透明可解释"),
        ("04", PURPLE, "运营可观测是持续优化的基础",
         "对话审计/业务标注/统计分析形成价值闭环"),
    ]

    num_w = px(50)
    gap = px(16)
    text_left = MARGIN_EMU + num_w + gap
    row_h = px(80)

    for i, (num, color, title, desc) in enumerate(conclusions):
        ry = y + i * (row_h + px(10))
        # 编号
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, MARGIN_EMU, ry + px(10), num_w, px(60))
        set_shape_fill(num_box, color)
        num_box.line.fill.background()
        add_textbox(slide, MARGIN_EMU, ry + px(24), num_w, px(32),
                    num, Pt(18), WHITE, True, PP_ALIGN.CENTER)

        # 标题
        add_textbox(slide, text_left, ry, SLIDE_W_EMU - text_left - MARGIN_EMU, px(24),
                    title, H4_SZ, TEXT_DARK, True)
        # 描述
        add_textbox(slide, text_left, ry + px(28), SLIDE_W_EMU - text_left - MARGIN_EMU, px(36),
                    desc, SMALL_SZ, TEXT_GRAY)


def build_p32(prs, layout):
    """P32 感谢页"""
    slide = prs.slides.add_slide(layout)
    add_cover_slide(slide,
                    "感谢聆听",
                    "AI 手机银行 11.0 — 技术方案设计",
                    "欢迎提问与交流",
                    "AIPmb 项目组 | 2026 年 6 月")


def build_divider(prs, layout, num, title):
    """章节分隔页"""
    slide = prs.slides.add_slide(layout)
    add_divider_slide(slide, num, title)
