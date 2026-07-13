"""主题常量：颜色、字体、布局尺寸、EMU 坐标系统"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 屏幕尺寸 (1440×820 px @ 96 DPI)
# ============================================================
PX_PER_INCH = 96
EMU_PER_INCH = 914400
PX_TO_EMU = EMU_PER_INCH // PX_PER_INCH  # 9525

SLIDE_W_PX = 1440
SLIDE_H_PX = 820
SLIDE_W_EMU = SLIDE_W_PX * PX_TO_EMU   # 13716000
SLIDE_H_EMU = SLIDE_H_PX * PX_TO_EMU   # 7810500

# ============================================================
# 间距系统 — 紧凑设计
# ============================================================
MARGIN = 32             # px
MARGIN_EMU = MARGIN * PX_TO_EMU       # 304800

HEADER_H = 64            # px, 标题栏高度
HEADER_H_EMU = HEADER_H * PX_TO_EMU
HEADER_GAP = 10          # px, 标题栏与内容区间距
HEADER_GAP_EMU = HEADER_GAP * PX_TO_EMU

BODY_TOP = MARGIN + HEADER_H + HEADER_GAP  # 106px
BODY_TOP_EMU = BODY_TOP * PX_TO_EMU

CONTENT_W = SLIDE_W_PX - 2 * MARGIN   # 1376
CONTENT_H = SLIDE_H_PX - BODY_TOP - MARGIN  # 682

# 居中辅助
CONTENT_W_EMU = CONTENT_W * PX_TO_EMU
SLIDE_CENTER_EMU = SLIDE_W_EMU // 2

# Logo 尺寸
LOGO_H = 52  # px
LOGO_H_EMU = LOGO_H * PX_TO_EMU
LOGO_W = 129  # px (等比缩放 258*52/104)
LOGO_W_EMU = LOGO_W * PX_TO_EMU
LOGO_GAP = 16

LOGO_PATH = "aipmb-ppt/yicheng.png"
LOGO_PATH_ABS = "/Users/pengjizhou/Documents/AIPmb/aipmb-ppt/yicheng.png"

# ============================================================
# 颜色系统 — 白色科技蓝主题
# ============================================================
ACCENT       = RGBColor(0x1A, 0x73, 0xE8)  # #1a73e8
ACCENT_DARK  = RGBColor(0x15, 0x57, 0xB0)  # #1557b0
ACCENT_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)  # #e8f0fe

WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x1A, 0x1A, 0x2E)   # #1a1a2e
DARK_SECONDARY = RGBColor(0x2D, 0x2D, 0x44)
TEXT_DARK    = RGBColor(0x21, 0x25, 0x29)   # #212529
TEXT_GRAY    = RGBColor(0x6C, 0x75, 0x7D)   # #6c757d
GRAY_50      = RGBColor(0xFA, 0xFB, 0xFC)
GRAY_100     = RGBColor(0xF8, 0xF9, 0xFA)
GRAY_200     = RGBColor(0xE9, 0xEC, 0xEF)

# 辅助色
BLUE         = RGBColor(0x25, 0x63, 0xEB)
GREEN        = RGBColor(0x05, 0x96, 0x69)
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
ORANGE       = RGBColor(0xEA, 0x58, 0x0C)

# 架构图层颜色
ARCH_BLUE_BG   = RGBColor(0xF0, 0xF9, 0xFF)
ARCH_CORE_BG   = RGBColor(0xEF, 0xF6, 0xFF)
ARCH_DOMAIN_BG = RGBColor(0xFF, 0xF7, 0xED)
ARCH_DATA_BG   = RGBColor(0xF0, 0xFD, 0xF4)
ARCH_INFRA_BG  = RGBColor(0xF5, 0xF3, 0xFF)

# 代码块配色
CODE_BG   = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT = RGBColor(0xE0, 0xE0, 0xE0)
CODE_COMMENT = RGBColor(0x6C, 0x70, 0x86)

# ============================================================
# 字体系统 — 放大字号
# ============================================================
FONT_FAMILY = 'PingFang SC'
FONT_FAMILY_CN = 'Microsoft YaHei'
CODE_FONT = 'SF Mono'

H1_SZ = Pt(48)   # 封面标题
H2_SZ = Pt(30)   # 页面标题
H3_SZ = Pt(22)   # 小标题
H4_SZ = Pt(18)   # 卡片标题
BODY_SZ = Pt(16) # 正文
SMALL_SZ = Pt(13) # 次要文字
CODE_SZ = Pt(11) # 代码
CAPTION_SZ = Pt(10) # 注释

# ============================================================
# 辅助函数
# ============================================================
def px(n):
    """像素 -> EMU"""
    return n * PX_TO_EMU


def center_x(width_emu):
    """给定宽度，返回居中 left 坐标"""
    return (SLIDE_W_EMU - width_emu) // 2


def set_font(run, size=BODY_SZ, color=TEXT_DARK, bold=False, family=None):
    """统一设置 run 的字体属性"""
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = family or FONT_FAMILY


def set_shape_fill(shape, color):
    """给 shape 设置纯色填充"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_shape_border(shape, color=None, width=None):
    """设置 shape 边框"""
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        if width:
            shape.line.width = Pt(width)


def add_textbox(slide, left_emu, top_emu, width_emu, height_emu,
                text="", font_size=BODY_SZ, color=TEXT_DARK, bold=False,
                alignment=PP_ALIGN.LEFT, font_family=None):
    """快捷添加文本框"""
    tb = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_family or FONT_FAMILY
    p.alignment = alignment
    return tb


def add_paragraph(tf, text, size=BODY_SZ, color=TEXT_DARK, bold=False, family=None):
    """在 text_frame 中追加段落"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = family or FONT_FAMILY
    p.space_before = Pt(4)
    p.space_after = Pt(0)
    return p


def add_rich_paragraph(tf, segments, alignment=PP_ALIGN.LEFT):
    """在 text_frame 中追加富文本段落"""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(4)
    p.space_after = Pt(0)
    for seg in segments:
        run = p.add_run()
        run.text = seg[0]
        run.font.size = seg[1] if len(seg) > 1 else BODY_SZ
        run.font.color.rgb = seg[2] if len(seg) > 2 else TEXT_DARK
        run.font.bold = seg[3] if len(seg) > 3 else False
        run.font.name = FONT_FAMILY
    return p
